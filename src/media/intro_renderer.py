"""
PowerPoint Intro Template Renderer

Renders a specific slide from a multi-slide PowerPoint template file.
Each slide is a complete intro design for different channels/purposes.

Usage:
    renderer = IntroTemplateRenderer("templates/intro_template.pptx")
    image = renderer.render(slide_index=0, title="News Title", article_image="path/to/image.jpg")
"""

import os
import subprocess
import tempfile
from pathlib import Path
from PIL import Image
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Emu
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class IntroTemplateRenderer:
    def __init__(self, template_path: str = "templates/intro_template.pptx"):
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx required: pip install python-pptx")
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        prs = Presentation(str(self.template_path))
        self.slide_count = len(prs.slides)
        # Get slide dimensions in pixels (assuming 96 DPI for screen)
        self.width = int(prs.slide_width.inches * 160)
        self.height = int(prs.slide_height.inches * 160)
    
    def list_slides(self) -> list:
        """List all slides with their index and any identifying text"""
        prs = Presentation(str(self.template_path))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = ''.join([p.text for p in shape.text_frame.paragraphs]).strip()
                    if text and "{{" not in text:
                        texts.append(text)
            slides.append({"index": i, "preview_text": texts[:3] if texts else []})
        return slides
    
    def find_slide_by_name(self, name: str) -> int:
        """Find slide index by matching name in slide text content"""
        slides = self.list_slides()
        name_lower = name.lower()
        
        for s in slides:
            for text in s.get('preview_text', []):
                if name_lower in text.lower():
                    return s['index']
        
        # Try as numeric index
        try:
            idx = int(name)
            if 0 <= idx < self.slide_count:
                return idx
        except ValueError:
            pass
        
        return 0  # Default to first slide

    def render(self, slide_index: int, title: str, article_image: str = None, 
               output_path: str = None) -> str:
        """
        Render a slide template with dynamic content.
        
        Args:
            slide_index: Which slide to use (0-based)
            title: Text to replace {{TITLE_HERE}} placeholder
            article_image: Path to article image (placed at bottom layer)
            output_path: Where to save the rendered image
            
        Returns:
            Path to rendered PNG image
        """
        if slide_index >= self.slide_count:
            raise ValueError(f"Slide {slide_index} not found. Template has {self.slide_count} slides.")
        
        # Load template
        prs = Presentation(str(self.template_path))
        slide = prs.slides[slide_index]
        
        # Replace title placeholder
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "{{TITLE_HERE}}" in run.text:
                            run.text = run.text.replace("{{TITLE_HERE}}", title)
        
        # Insert article image at bottom layer if provided
        if article_image and os.path.exists(article_image):
            self._insert_background_image(slide, article_image, prs)
        
        # Save modified presentation to temp file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs.save(tmp.name)
            tmp_pptx = tmp.name
        
        # Determine output path
        if output_path is None:
            output_path = tempfile.mktemp(suffix=".png")
        
        try:
            return self._convert_to_image(tmp_pptx, output_path)
        finally:
            os.unlink(tmp_pptx)
    
    def _insert_background_image(self, slide, image_path: str, prs):
        """Insert image at the bottom layer (behind all other elements)"""
        from pptx.util import Inches
        
        # Add picture at full slide size
        pic = slide.shapes.add_picture(
            image_path, 
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )
        
        # Move to back (bottom of z-order)
        spTree = slide.shapes._spTree
        pic_element = pic._element
        spTree.remove(pic_element)
        spTree.insert(2, pic_element)  # Insert after background, before other shapes

    def _convert_to_image(self, pptx_path: str, output_path: str) -> str:
        """Convert PowerPoint slide to image using LibreOffice"""
        output_dir = os.path.dirname(output_path) or "."
        os.makedirs(output_dir, exist_ok=True)
        
        try:
            # Use LibreOffice to convert - preserves all formatting, transparency, etc.
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "png", "--outdir", output_dir, pptx_path],
                capture_output=True, text=True, timeout=120
            )
            
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                generated_png = os.path.join(output_dir, f"{base_name}.png")
                
                if os.path.exists(generated_png):
                    # Resize to exact dimensions
                    img = Image.open(generated_png)
                    img = img.resize((self.width, self.height), Image.Resampling.LANCZOS)
                    img.save(output_path)
                    
                    if generated_png != output_path:
                        os.unlink(generated_png)
                    
                    return output_path
                    
        except subprocess.TimeoutExpired:
            print("LibreOffice conversion timed out")
        except FileNotFoundError:
            print("LibreOffice not found - install with: sudo apt install libreoffice")
        except Exception as e:
            print(f"LibreOffice conversion failed: {e}")
        
        raise RuntimeError("Failed to convert PowerPoint to image. Ensure LibreOffice is installed.")
    
    def render_to_numpy(self, slide_index: int, title: str, article_image: str = None) -> np.ndarray:
        """Render directly to numpy array for MoviePy integration"""
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp_path = tmp.name
        
        try:
            self.render(slide_index, title, article_image, tmp_path)
            return np.array(Image.open(tmp_path).convert('RGB'))
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)


# Convenience function
def render_intro(slide_index: int, title: str, article_image: str = None,
                 template_path: str = "templates/intro_template.pptx",
                 output_path: str = None) -> str:
    """Quick render function"""
    renderer = IntroTemplateRenderer(template_path)
    return renderer.render(slide_index, title, article_image, output_path)


if __name__ == "__main__":
    # Demo
    renderer = IntroTemplateRenderer()
    print(f"Template: {renderer.template_path}")
    print(f"Slides: {renderer.slide_count}")
    print(f"Resolution: {renderer.width}x{renderer.height}")
    print("\nSlides:")
    for s in renderer.list_slides():
        print(f"  [{s['index']}] {s['preview_text']}")
