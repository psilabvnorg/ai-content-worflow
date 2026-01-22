"""
Media Module - Consolidated TTS, subtitle, and video generation.

This module handles all media generation: voice synthesis, subtitle creation,
and video composition with effects.
"""
import os
import wave
import subprocess
import tempfile
from pathlib import Path
import numpy as np
import pysrt
import whisper
import torch
import soundfile as sf
from PIL import Image, ImageFilter, ImageDraw, ImageFont

# MoviePy 2.x compatible imports
try:
    from moviepy import VideoClip, ImageClip, AudioFileClip, CompositeVideoClip, concatenate_videoclips
    from moviepy import VideoFileClip, AudioClip, concatenate_audioclips, CompositeAudioClip
    import moviepy.video.fx as vfx
    import moviepy.audio.fx as afx
    MOVIEPY_VERSION = 2
except ImportError:
    # MoviePy 1.x fallback
    from moviepy.video.VideoClip import VideoClip, ImageClip, AudioClip
    from moviepy.audio.io.AudioFileClip import AudioFileClip
    from moviepy.video.compositing.CompositeVideoClip import CompositeVideoClip
    from moviepy.video.compositing.concatenate import concatenate_videoclips
    from moviepy.video.io.VideoFileClip import VideoFileClip
    from moviepy.audio.AudioClip import concatenate_audioclips, CompositeAudioClip
    import moviepy.video.fx.all as vfx
    import moviepy.audio.fx.all as afx
    MOVIEPY_VERSION = 1

try:
    from pptx import Presentation
    from pptx.util import Emu
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class MediaGenerator:
    '''
    Unified media generation class for TTS, subtitles, and video composition.
    
    Responsibilities:
    - Generate Vietnamese voice-over using VieNeu-TTS
    - Create synchronized subtitles with Whisper
    - Compose final video with effects and overlays
    - Render PowerPoint intro templates
    '''
    
    def __init__(self, voice: str = "binh", resolution=(1080, 1920), fps=30):
        '''
        Initialize media generator with voice and video settings.
        
        Args:
            voice: Voice name for TTS (binh, tuyen, nguyen, etc.)
            resolution: Video resolution (width, height)
            fps: Frames per second
        '''
        self.voice_name = voice
        self.width, self.height = resolution
        self.fps = fps
        self.tts = None
        self.current_voice = None
        self.whisper_model = None
        self._init_tts()
        print(f"✓ MediaGenerator initialized (Voice: {voice}, Resolution: {resolution[0]}x{resolution[1]})")
    
    def _init_tts(self):
        '''Initialize VieNeu-TTS for Vietnamese speech synthesis.'''
        try:
            from vieneu import Vieneu
            has_cuda = torch.cuda.is_available()
            device = "cuda" if has_cuda else "cpu"
            local_path = "models/VieNeu-TTS"
            
            if has_cuda and os.path.exists(local_path):
                self.tts = Vieneu(backbone_repo=local_path, backbone_device=device, codec_device=device)
            elif has_cuda:
                self.tts = Vieneu(backbone_repo="pnnbao-ump/VieNeu-TTS-0.3B", backbone_device=device, codec_device=device)
            else:
                self.tts = Vieneu(backbone_repo="pnnbao-ump/VieNeu-TTS-0.3B-q8-gguf")
            
            voices = self.tts.list_preset_voices()
            available = [v[1] if isinstance(v, tuple) else v for v in voices]
            voice_map = {"binh": "Binh", "tuyen": "Tuyen", "nguyen": "Nguyen", "son": "Son", 
                        "vinh": "Vinh", "huong": "Huong", "ly": "Ly", "ngoc": "Ngoc", 
                        "doan": "Doan", "dung": "Dung"}
            target = voice_map.get(self.voice_name.lower(), self.voice_name.capitalize())
            
            if target in available:
                self.current_voice = self.tts.get_preset_voice(target)
            else:
                for v in ["Binh", "Tuyen"]:
                    if v in available:
                        self.current_voice = self.tts.get_preset_voice(v)
                        break
            
            print(f"   ✓ VieNeu-TTS ready ({device.upper()} mode)")
        except Exception as e:
            print(f"   ✗ VieNeu-TTS init failed: {e}")
            self.tts = None
    
    def generate_audio(self, text: str, output_path: str) -> str:
        '''
        Generate speech audio from text.
        
        Args:
            text: Text to synthesize
            output_path: Path to save audio file
            
        Returns:
            Path to generated audio file
        '''
        if not self.tts:
            raise RuntimeError("VieNeu-TTS not initialized")
        
        clean_text = text.replace("... ", ". ").replace(" ... ", ". ")
        audio = self.tts.infer(text=clean_text, voice=self.current_voice, temperature=0.8, top_k=50)
        
        if output_path.endswith('.mp3'):
            temp_wav = output_path.replace('.mp3', '_temp.wav')
            sf.write(temp_wav, audio, 24000)
            try:
                subprocess.run(['ffmpeg', '-i', temp_wav, '-codec:a', 'libmp3lame', '-qscale:a', '2', '-y', output_path],
                             capture_output=True, check=True)
                os.remove(temp_wav)
            except:
                os.rename(temp_wav, output_path)
        else:
            sf.write(output_path, audio, 24000)
        
        print(f"✓ Audio generated: {output_path}")
        return output_path
    
    def get_audio_duration(self, audio_path: str) -> float:
        '''Get audio duration in seconds.'''
        try:
            data, rate = sf.read(audio_path)
            return len(data) / rate
        except:
            result = subprocess.run(['ffprobe', '-v', 'error', '-show_entries', 'format=duration',
                                   '-of', 'default=noprint_wrappers=1:nokey=1', audio_path],
                                  capture_output=True, text=True)
            return float(result.stdout.strip())
    
    def generate_subtitles(self, audio_path: str, output_path: str, original_script: str = None) -> str:
        '''
        Generate synchronized subtitles using Whisper.
        
        Args:
            audio_path: Path to audio file
            output_path: Path to save SRT file
            original_script: Original corrected script for alignment
            
        Returns:
            Path to generated subtitle file
        '''
        try:
            if torch.cuda.is_available():
                torch.cuda.empty_cache()
            
            if not self.whisper_model:
                for model_name in ["base", "small", "tiny"]:
                    try:
                        self.whisper_model = whisper.load_model(model_name)
                        break
                    except torch.cuda.OutOfMemoryError:
                        torch.cuda.empty_cache()
                        continue
                if not self.whisper_model:
                    self.whisper_model = whisper.load_model("base", device="cpu")
            
            result = self.whisper_model.transcribe(audio_path, language="vi", word_timestamps=True, 
                                                   verbose=False, fp16=torch.cuda.is_available())
            
            whisper_words = []
            for segment in result['segments']:
                if 'words' in segment:
                    for w in segment['words']:
                        whisper_words.append({'word': w['word'].strip(), 'start': w['start'], 'end': w['end']})
            
            if not whisper_words:
                return self._fallback_subtitles(audio_path, output_path, original_script)
            
            corrected_words = (original_script if original_script else result['text']).split()
            aligned = self._align_words(whisper_words, corrected_words)
            
            subs = pysrt.SubRipFile()
            idx, phrase_idx = 0, 1
            while idx < len(aligned):
                phrase_words, start, end = [], aligned[idx]['start'], aligned[idx]['end']
                for j in range(4):
                    if idx + j >= len(aligned):
                        break
                    pair = aligned[idx + j]
                    phrase_words.append(pair['corrected_word'])
                    end = pair['end']
                    if pair['corrected_word'].endswith(('.', '!', '?', ':')):
                        break
                if phrase_words:
                    subs.append(self._create_sub_item(phrase_idx, start, end, ' '.join(phrase_words)))
                    phrase_idx += 1
                idx += len(phrase_words)
            
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            subs.save(output_path, encoding='utf-8')
            print(f"✓ Generated {len(subs)} subtitle phrases")
            return output_path
        except Exception as e:
            print(f"Subtitle generation error: {e}")
            return self._fallback_subtitles(audio_path, output_path, original_script)
    
    def _align_words(self, whisper_words: list, corrected_words: list) -> list:
        '''Align corrected words with Whisper timing.'''
        if len(whisper_words) == len(corrected_words):
            return [{'corrected_word': corrected_words[i], 'start': whisper_words[i]['start'], 
                    'end': whisper_words[i]['end']} for i in range(len(corrected_words))]
        
        aligned, whisper_idx = [], 0
        total_duration = whisper_words[-1]['end'] - whisper_words[0]['start']
        time_per_word = total_duration / len(corrected_words)
        
        for i, corrected_word in enumerate(corrected_words):
            best_idx, best_score = whisper_idx, 0
            for j in range(whisper_idx, min(whisper_idx + 3, len(whisper_words))):
                w_word = whisper_words[j]['word'].lower()
                c_word = corrected_word.lower().strip('.,!?:;')
                score = 1.0 if w_word == c_word else (0.7 if w_word.startswith(c_word[:3]) or c_word.startswith(w_word[:3]) else 0.0)
                if score > best_score:
                    best_score, best_idx = score, j
            
            if best_score > 0.5 and best_idx < len(whisper_words):
                aligned.append({'corrected_word': corrected_word, 'start': whisper_words[best_idx]['start'], 
                              'end': whisper_words[best_idx]['end']})
                whisper_idx = best_idx + 1
            else:
                start = whisper_words[0]['start'] + (i * time_per_word)
                aligned.append({'corrected_word': corrected_word, 'start': start, 'end': start + time_per_word})
        return aligned
    
    def _create_sub_item(self, index: int, start: float, end: float, text: str) -> pysrt.SubRipItem:
        '''Create SubRipItem with proper timing.'''
        return pysrt.SubRipItem(
            index=index,
            start={'hours': int(start//3600), 'minutes': int((start%3600)//60), 
                  'seconds': int(start%60), 'milliseconds': int((start%1)*1000)},
            end={'hours': int(end//3600), 'minutes': int((end%3600)//60), 
                'seconds': int(end%60), 'milliseconds': int((end%1)*1000)},
            text=text
        )
    
    def _fallback_subtitles(self, audio_path: str, output_path: str, script: str) -> str:
        '''Fallback subtitle generation.'''
        try:
            data, rate = sf.read(audio_path)
            duration = len(data) / rate
            text = script if script else 'Subtitle generation failed'
            words = text.split()
            chunks = [' '.join(words[i:i+10]) for i in range(0, len(words), 10)]
            time_per_chunk = duration / len(chunks) if chunks else duration
            
            subs = pysrt.SubRipFile()
            for i, chunk in enumerate(chunks):
                subs.append(self._create_sub_item(i+1, i*time_per_chunk, (i+1)*time_per_chunk, chunk))
            
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            subs.save(output_path, encoding='utf-8')
            return output_path
        except:
            subs = pysrt.SubRipFile()
            subs.append(pysrt.SubRipItem(index=1, start={'hours':0,'minutes':0,'seconds':0,'milliseconds':0},
                                        end={'hours':0,'minutes':1,'seconds':0,'milliseconds':0},
                                        text="[Subtitle generation failed]"))
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            subs.save(output_path, encoding='utf-8')
            return output_path
    
    def compose_video(self, images: list, audio_path: str, subtitle_path: str, output_path: str,
                     audio_duration: float, title: str = None, background_music: str = None,
                     typing_sfx: str = None, broll_videos: list = None, template: str = None,
                     intro_duration: float = 3.0) -> str:
        '''
        Compose final TikTok video with effects.
        
        Args:
            images: List of image paths
            audio_path: Path to voice-over audio
            subtitle_path: Path to SRT subtitle file
            output_path: Path to save final video
            audio_duration: Duration of voice-over
            title: Video title for intro
            background_music: Path to background music
            typing_sfx: Path to typing sound effect
            broll_videos: List of B-roll video paths
            template: PowerPoint template slide name/index
            intro_duration: Intro duration (None = full video)
            
        Returns:
            Path to generated video
        '''
        if not images:
            raise ValueError("No images provided")
        
        full_video_intro = intro_duration is None
        actual_intro_duration = audio_duration if full_video_intro else intro_duration
        has_separate_intro = not full_video_intro and title
        
        total_media = len(images) + (len(broll_videos) if broll_videos else 0)
        duration_per_media = audio_duration / total_media
        
        clips, intro_overlay = [], None
        
        if title and full_video_intro:
            intro_overlay = self._create_intro_overlay(title, images[0], audio_duration, template)
        elif title and has_separate_intro:
            intro_clip = self._create_intro_clip(title, images[0], actual_intro_duration, template)
            if MOVIEPY_VERSION == 2:
                intro_clip = intro_clip.with_effects([vfx.CrossFadeOut(0.5)])
            else:
                intro_clip = intro_clip.crossfadeout(0.5)
            clips.append(intro_clip)
        
        for img_path in images:
            clips.append(self._create_effect_clip(img_path, duration_per_media))
        
        if broll_videos:
            for video_path in broll_videos:
                broll = self._create_broll_clip(video_path, duration_per_media)
                if broll:
                    clips.append(broll)
        
        video = concatenate_videoclips(clips, method="compose")
        if intro_overlay:
            video = CompositeVideoClip([video, intro_overlay])
        
        voice_audio = AudioFileClip(audio_path)
        if has_separate_intro:
            silence = AudioClip(lambda t: 0, duration=actual_intro_duration, fps=44100)
            voice_audio = concatenate_audioclips([silence, voice_audio])
        
        audio_layers = [voice_audio]
        total_duration = (actual_intro_duration if has_separate_intro else 0) + audio_duration
        
        if title and actual_intro_duration > 0:
            typing = self._add_typing_sfx(actual_intro_duration, typing_sfx)
            if typing:
                audio_layers.append(typing)
        
        if background_music and os.path.exists(background_music):
            try:
                bg = AudioFileClip(background_music)
                if bg.duration < total_duration:
                    loops = int(total_duration / bg.duration) + 1
                    bg = concatenate_audioclips([bg] * loops)
                bg = bg.subclipped(0, total_duration) if MOVIEPY_VERSION == 2 else bg.subclip(0, total_duration)
                if MOVIEPY_VERSION == 2:
                    bg = bg.with_effects([afx.MultiplyVolume(0.15)])
                else:
                    bg = bg.volumex(0.15)
                audio_layers.append(bg)
            except Exception as e:
                print(f"Background music error: {e}")
        video = video.with_audio(CompositeAudioClip(audio_layers)) if MOVIEPY_VERSION == 2 else video.set_audio(CompositeAudioClip(audio_layers))
        
        time_offset = actual_intro_duration if has_separate_intro else 0
        video = self._add_subtitles(video, subtitle_path, time_offset)
        
        video.write_videofile(output_path, fps=self.fps, codec='libx264', audio_codec='aac',
                            bitrate='6000k', preset='medium', threads=4)
        print(f"✓ Video created: {output_path}")
        return output_path
    
    def _create_effect_clip(self, image_path: str, duration: float) -> VideoClip:
        '''Create clip with blurred background and pan effect.'''
        img = Image.open(image_path).convert('RGB')
        img_ratio = img.width / img.height
        
        bg = img.resize((self.width, self.height), Image.Resampling.LANCZOS).filter(ImageFilter.GaussianBlur(40))
        bg_array = np.array(bg)
        
        base_width = self.width
        base_height = int(base_width / img_ratio)
        if base_height > self.height * 0.7:
            base_height = int(self.height * 0.7)
            base_width = int(base_height * img_ratio)
        if base_height < self.height * 0.5:
            base_height = int(self.height * 0.5)
            base_width = int(base_height * img_ratio)
        
        pan_width, pan_height = int(base_width * 1.15), int(base_height * 1.15)
        
        def make_frame(t):
            progress = t / duration
            eased = progress * progress * (3 - 2 * progress)
            frame = bg_array.copy()
            current = np.array(img.resize((pan_width, pan_height), Image.Resampling.LANCZOS))
            max_pan = max(0, pan_width - self.width)
            x, y = -int(max_pan * eased), (self.height - pan_height) // 2
            sx1, sy1 = max(0, -x), max(0, -y)
            sx2, sy2 = min(pan_width, self.width - x), min(pan_height, self.height - y)
            dx1, dy1 = max(0, x), max(0, y)
            dx2, dy2 = min(self.width, x + pan_width), min(self.height, y + pan_height)
            if sx2 > sx1 and sy2 > sy1:
                frame[dy1:dy2, dx1:dx2] = current[sy1:sy2, sx1:sx2]
            return frame
        
        return VideoClip(make_frame, duration=duration)
    
    def _create_broll_clip(self, video_path: str, target_duration: float) -> VideoClip:
        '''Create B-roll video clip with effects.'''
        try:
            if not os.path.exists(video_path):
                return None
            broll = VideoFileClip(video_path).without_audio()
            if broll.duration > target_duration:
                broll = broll.subclipped(0, target_duration) if MOVIEPY_VERSION == 2 else broll.subclip(0, target_duration)
            elif broll.duration < target_duration:
                loops = int(target_duration / broll.duration) + 1
                broll = concatenate_videoclips([broll] * loops)
                broll = broll.subclipped(0, target_duration) if MOVIEPY_VERSION == 2 else broll.subclip(0, target_duration)
            return self._resize_broll(broll)
        except Exception as e:
            print(f"B-roll error: {e}")
            return None
    
    def _resize_broll(self, video_clip: VideoClip) -> VideoClip:
        '''Resize B-roll with blurred background and pan.'''
        ow, oh = video_clip.size
        ratio = ow / oh
        nw, nh = self.width, int(self.width / ratio)
        if nh > self.height * 0.7:
            nh, nw = int(self.height * 0.7), int(nh * ratio)
        if nh < self.height * 0.5:
            nh, nw = int(self.height * 0.5), int(nh * ratio)
        pan_w, pan_h = int(nw * 1.15), int(nh * 1.15)
        resized = video_clip.resized((pan_w, pan_h))
        bg = Image.fromarray(video_clip.get_frame(0)).resize((self.width, self.height), Image.Resampling.LANCZOS)
        bg_array = np.array(bg.filter(ImageFilter.GaussianBlur(40)))
        duration = video_clip.duration
        
        def make_frame(t):
            progress, eased = t / duration, (t/duration)**2 * (3 - 2*(t/duration))
            frame = bg_array.copy()
            vf = resized.get_frame(t)
            max_pan = max(0, pan_w - self.width)
            x, y = -int(max_pan * eased), (self.height - pan_h) // 2
            sx1, sy1 = max(0, -x), max(0, -y)
            sx2, sy2 = min(pan_w, self.width - x), min(pan_h, self.height - y)
            dx1, dy1 = max(0, x), max(0, y)
            dx2, dy2 = min(self.width, x + pan_w), min(self.height, y + pan_h)
            if sx2 > sx1 and sy2 > sy1:
                frame[dy1:dy2, dx1:dx2] = vf[sy1:sy2, sx1:sx2]
            return frame
        
        return VideoClip(make_frame, duration=duration)
    
    def _add_typing_sfx(self, duration: float, custom_path: str = None) -> AudioFileClip:
        '''Add typing sound effect.'''
        paths = [custom_path, 'assets/typing.mp3', 'assets/typing.wav']
        sfx_path = next((p for p in paths if p and os.path.exists(p)), None)
        if not sfx_path:
            return None
        try:
            audio = AudioFileClip(sfx_path)
            if audio.duration < duration:
                audio = concatenate_audioclips([audio] * (int(duration / audio.duration) + 1))
            audio = audio.subclipped(0, duration) if MOVIEPY_VERSION == 2 else audio.subclip(0, duration)
            if MOVIEPY_VERSION == 2:
                return audio.with_effects([afx.MultiplyVolume(0.3)])
            else:
                return audio.volumex(0.3)
        except:
            return None
    
    def _add_subtitles(self, video: CompositeVideoClip, subtitle_path: str, time_offset: float = 0.0) -> CompositeVideoClip:
        '''Add subtitles to video.'''
        try:
            subs = pysrt.open(subtitle_path, encoding='utf-8')
            clips = []
            for sub in subs:
                start = (sub.start.ordinal / 1000) + time_offset
                duration = (sub.end.ordinal / 1000) + time_offset - start
                img = self._create_subtitle_image(sub.text)
                clip = ImageClip(img, duration=duration).with_position(('center', self.height - 250)).with_start(start)
                clips.append(clip)
            video = CompositeVideoClip([video] + clips)
            print(f"   ✓ Added {len(clips)} subtitle clips")
        except Exception as e:
            print(f"Subtitle overlay error: {e}")
        return video
    
    def _create_subtitle_image(self, text: str, max_width: int = 980) -> np.ndarray:
        '''Create subtitle image with text.'''
        font = self._get_font(38)
        temp = Image.new('RGB', (max_width, 500), 'black')
        draw = ImageDraw.Draw(temp)
        words, lines, cur = text.split(), [], []
        for w in words:
            test = ' '.join(cur + [w])
            if draw.textbbox((0, 0), test, font=font)[2] <= max_width - 40:
                cur.append(w)
            else:
                if cur:
                    lines.append(' '.join(cur))
                cur = [w]
        if cur:
            lines.append(' '.join(cur))
        line_h, total_h = 48, len(lines) * 48 + 30
        img = Image.new('RGBA', (max_width, total_h), (0, 0, 0, 220))
        draw = ImageDraw.Draw(img)
        y = 15
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=font)
            x = (max_width - (bbox[2] - bbox[0])) // 2
            color = '#FFFF99' if any(p in line for p in '.!?:') else 'white'
            draw.text((x, y), line, font=font, fill=color)
            y += line_h
        return np.array(img)
    
    def _get_font(self, size: int):
        '''Get font with fallback.'''
        for p in ['/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
                  '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf']:
            if os.path.exists(p):
                return ImageFont.truetype(p, size)
        return ImageFont.load_default()
    
    def _create_intro_clip(self, title: str, image_path: str, duration: float, template: str = None) -> VideoClip:
        '''Create intro clip using PowerPoint template or fallback.'''
        if template and PPTX_AVAILABLE:
            try:
                intro_array = self._render_pptx_template(template, title, image_path)
                return ImageClip(intro_array, duration=duration)
            except Exception as e:
                print(f"PowerPoint template error: {e}")
        return self._create_fallback_intro(title, image_path, duration)
    
    def _create_intro_overlay(self, title: str, image_path: str, duration: float, template: str = None) -> VideoClip:
        '''Create intro overlay for full video duration.'''
        if template and PPTX_AVAILABLE:
            try:
                intro_array = self._render_pptx_template(template, title, image_path)
                return ImageClip(intro_array, duration=duration).with_position(('center', 'center'))
            except Exception as e:
                print(f"PowerPoint template error: {e}")
        return self._create_fallback_intro_overlay(title, image_path, duration)
    
    def _render_pptx_template(self, template: str, title: str, image_path: str) -> np.ndarray:
        '''Render PowerPoint template to numpy array.'''
        prs = Presentation("templates/intro_template.pptx")
        slide_idx = 0
        try:
            slide_idx = int(template)
        except:
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = ''.join([p.text for p in shape.text_frame.paragraphs]).strip()
                        if template.lower() in text.lower():
                            slide_idx = i
                            break
        
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "{{TITLE_HERE}}" in run.text:
                            run.text = run.text.replace("{{TITLE_HERE}}", title)
        
        if image_path and os.path.exists(image_path):
            pic = slide.shapes.add_picture(image_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
            spTree = slide.shapes._spTree
            spTree.remove(pic._element)
            spTree.insert(2, pic._element)
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            prs.save(tmp.name)
            tmp_pptx = tmp.name
        
        try:
            output_dir = tempfile.mkdtemp()
            subprocess.run(["libreoffice", "--headless", "--convert-to", "png", "--outdir", output_dir, tmp_pptx],
                         capture_output=True, timeout=120, check=True)
            base_name = os.path.splitext(os.path.basename(tmp_pptx))[0]
            png_path = os.path.join(output_dir, f"{base_name}.png")
            img = Image.open(png_path).resize((self.width, self.height), Image.Resampling.LANCZOS)
            result = np.array(img.convert('RGB'))
            os.unlink(png_path)
            os.rmdir(output_dir)
            return result
        finally:
            os.unlink(tmp_pptx)
    
    def _create_fallback_intro(self, title: str, image_path: str, duration: float) -> VideoClip:
        '''Fallback intro without PowerPoint.'''
        img = Image.open(image_path).convert('RGB')
        img = self._resize_image_full(img)
        if os.path.exists('assets/tiktok_background.png'):
            overlay = Image.open('assets/tiktok_background.png').convert('RGBA')
            overlay = overlay.resize((self.width, self.height), Image.Resampling.LANCZOS)
            img = Image.alpha_composite(img.convert('RGBA'), overlay)
        else:
            img = img.convert('RGBA')
        draw = ImageDraw.Draw(img)
        font = self._get_font(48)
        words, lines, cur = title.split(), [], []
        for w in words:
            test = ' '.join(cur + [w])
            if draw.textbbox((0, 0), test, font=font)[2] <= self.width - 120:
                cur.append(w)
            else:
                if cur:
                    lines.append(' '.join(cur))
                cur = [w]
        if cur:
            lines.append(' '.join(cur))
        y = self.height // 2 + 200
        for line in lines:
            draw.text((60, y), line, font=font, fill='white')
            y += 60
        return ImageClip(np.array(img.convert('RGB')), duration=duration)
    
    def _create_fallback_intro_overlay(self, title: str, image_path: str, duration: float) -> VideoClip:
        '''Fallback intro overlay with semi-transparent background.'''
        img = Image.new('RGBA', (self.width, self.height), (0, 0, 0, 0))
        overlay_bg = Image.new('RGBA', (self.width, 300), (0, 0, 0, 180))
        img.paste(overlay_bg, (0, 0))
        draw = ImageDraw.Draw(img)
        font = self._get_font(48)
        words, lines, cur = title.split(), [], []
        for w in words:
            test = ' '.join(cur + [w])
            if draw.textbbox((0, 0), test, font=font)[2] <= self.width - 120:
                cur.append(w)
            else:
                if cur:
                    lines.append(' '.join(cur))
                cur = [w]
        if cur:
            lines.append(' '.join(cur))
        y = 60
        for line in lines:
            draw.text((60, y), line, font=font, fill='white')
            y += 60
        return ImageClip(np.array(img), duration=duration).with_position(('center', 'top'))
    
    def _resize_image_full(self, img: Image.Image) -> Image.Image:
        '''Resize image to fill screen.'''
        ratio, target = img.width / img.height, self.width / self.height
        if ratio > target:
            nh, nw = self.height, int(self.height * ratio)
            img = img.resize((nw, nh), Image.Resampling.LANCZOS)
            left = (nw - self.width) // 2
            img = img.crop((left, 0, left + self.width, self.height))
        else:
            nw, nh = self.width, int(self.width / ratio)
            img = img.resize((nw, nh), Image.Resampling.LANCZOS)
            top = (nh - self.height) // 2
            img = img.crop((0, top, self.width, top + self.height))
        return img
